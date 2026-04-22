from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "deck-assets"
OUT = ROOT / "ReliefConnect-Demo-Flow.pptx"


def add_title_band(slide, title: str, subtitle: str = ""):
    band = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(0.95)
    )  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(44, 69, 112)
    band.line.fill.background()

    t = band.text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 250, 246)

    if subtitle:
        p2 = t.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(220, 232, 248)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(92, 103, 117)
    p.alignment = PP_ALIGN.RIGHT


def add_full_image_slide(prs, title, subtitle, image_name, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, title, subtitle)
    img = ASSETS / image_name
    slide.shapes.add_picture(str(img), Inches(0.4), Inches(1.15), Inches(8.8), Inches(5.7))

    box = slide.shapes.add_textbox(Inches(9.35), Inches(1.2), Inches(3.6), Inches(5.7))
    tf = box.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(49, 58, 70)
    add_footer(slide, "ReliefConnect Hackathon Demo")


def add_side_by_side_slide(prs, title, subtitle, left_image, right_image, left_label, right_label, flow):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, title, subtitle)

    slide.shapes.add_picture(str(ASSETS / left_image), Inches(0.45), Inches(1.35), Inches(6.15), Inches(4.9))
    slide.shapes.add_picture(str(ASSETS / right_image), Inches(6.75), Inches(1.35), Inches(6.15), Inches(4.9))

    lbox = slide.shapes.add_textbox(Inches(0.45), Inches(6.3), Inches(6.15), Inches(0.35))
    ltf = lbox.text_frame
    ltf.text = f"Origin: {left_label}"
    ltf.paragraphs[0].font.size = Pt(12)
    ltf.paragraphs[0].font.bold = True

    rbox = slide.shapes.add_textbox(Inches(6.75), Inches(6.3), Inches(6.15), Inches(0.35))
    rtf = rbox.text_frame
    rtf.text = f"Dependent: {right_label}"
    rtf.paragraphs[0].font.size = Pt(12)
    rtf.paragraphs[0].font.bold = True

    flow_box = slide.shapes.add_textbox(Inches(0.45), Inches(6.65), Inches(12.45), Inches(0.45))
    ftf = flow_box.text_frame
    ftf.clear()
    p = ftf.paragraphs[0]
    p.text = f"Flow: {flow}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 69, 112)
    add_footer(slide, "ReliefConnect dependency mapping")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "ReliefConnect", "Feeding communities when it matters most")
    slide.shapes.add_picture(str(ASSETS / "01-landing.png"), Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.6))
    add_footer(slide, "Demo story: resident -> hub -> admin -> resident loop")

    # Slide 2
    add_full_image_slide(
        prs,
        "Resident Onboarding & Hub Discovery",
        "Start-to-end flow begins with household registration and live hub view",
        "02-user-register.png",
        [
            "Residents register and sign in with household profile.",
            "Hub map and lane health guide where to act first.",
            "Disaster clock controls phase behavior for demo.",
        ],
    )

    # Slide 3
    add_full_image_slide(
        prs,
        "Reservation + Shelter Coordination",
        "Hub intake transitions based on disaster timing and shelter coordination",
        "04-user-warehouse.png",
        [
            "Residents reserve balanced baskets by nutrition lane.",
            "At crisis window, hub intake closes and shelters become primary.",
            "Rule-based messaging keeps residents informed in-app.",
        ],
    )

    # Slide 4
    add_side_by_side_slide(
        prs,
        "Dependency A: Donations Flow",
        "Cross-role dependency (Resident action -> Hub processing)",
        "06-user-donate.png",
        "07-hub-donations.png",
        "Resident creates donation ticket",
        "Hub accepts/closes donation ticket",
        "User donates -> Hub verifies and closes -> Inventory and points pipeline updates",
    )

    # Slide 5
    add_side_by_side_slide(
        prs,
        "Dependency B: Operational Feedback Loop",
        "Hub requests + Admin observability + execution feedback",
        "08-hub-dashboard.png",
        "09-admin-dashboard.png",
        "Hub operations and requests",
        "Master admin regional observability",
        "Hub submits/operates -> Admin monitors regions and allocates decisions -> System state updates",
    )

    # Slide 6
    add_side_by_side_slide(
        prs,
        "Dependency C: Impact Back to Resident",
        "Closed-loop value delivery to the same resident journey",
        "07-hub-donations.png",
        "11-user-redeem.png",
        "Hub completion of donation lifecycle",
        "Resident reward redemption",
        "Hub closes verified ticket -> Credits posted -> Resident redeems by nutrition need",
    )

    # Slide 7
    add_full_image_slide(
        prs,
        "Continuity & Accountability",
        "Real-time holds, history, and operational transparency",
        "10-user-reservations.png",
        [
            "Reservation timeline provides countdown and status visibility.",
            "Multi-role state is synchronized across resident, hub, and admin.",
            "Designed for demo clarity while preserving dependency logic.",
        ],
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "Demo Recap", "End-to-end ReliefConnect functionality")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(12.0), Inches(5.5))
    tf = box.text_frame
    tf.clear()
    points = [
        "1) Resident registers, discovers hub status, and reserves essentials.",
        "2) Resident donations are dependency-linked to hub acceptance and closure.",
        "3) Hub/admin collaboration drives regional metrics and operational decisions.",
        "4) Verified actions return value to residents via credits and redeem flows.",
        "5) Shelter and disaster phase controls enforce safe, policy-aware behavior.",
        "Demo credentials: admin@warehouse.com / 123456",
    ]
    for i, item in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22 if i == 0 else 18)
        p.font.color.rgb = RGBColor(44, 69, 112) if i == 0 else RGBColor(49, 58, 70)
        p.space_after = Pt(9)
    add_footer(slide, "ReliefConnect-Demo-Flow.pptx")

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
